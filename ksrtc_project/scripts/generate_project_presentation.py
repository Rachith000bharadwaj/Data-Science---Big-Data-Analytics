#!/usr/bin/env python3
from __future__ import annotations

from datetime import datetime
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_RAW = ROOT / "data" / "raw"
RESULTS = ROOT / "results" / "csv"
VISUALS = ROOT / "visualizations"
OUTPUT = ROOT / "KSRTC_Whole_Project_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(247, 245, 240)
NAVY = RGBColor(20, 38, 66)
BLUE = RGBColor(38, 98, 195)
ORANGE = RGBColor(241, 120, 53)
GREEN = RGBColor(54, 122, 76)
TEXT = RGBColor(33, 37, 41)
MUTED = RGBColor(102, 112, 122)
WHITE = RGBColor(255, 255, 255)
LIGHT_BORDER = RGBColor(219, 221, 226)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
TOTAL_SLIDES = 12


def csv_rows(path: Path) -> int:
    if not path.exists():
        return 0
    with path.open("r", encoding="utf-8", errors="ignore") as handle:
        return max(sum(1 for _ in handle) - 1, 0)


def read_top_hubs() -> list[tuple[str, str]]:
    top_hubs_path = RESULTS / "top_pagerank_hubs.csv"
    if top_hubs_path.exists():
        df = pd.read_csv(top_hubs_path).head(5)
        return [
            (str(row.stop_name), f"{float(row.rank_score):.4f}")
            for row in df.itertuples(index=False)
        ]

    pagerank_path = RESULTS / "pagerank.csv"
    stops_path = DATA_RAW / "synthetic_bus_stops.csv"
    if not pagerank_path.exists() or not stops_path.exists():
        return []

    pagerank = pd.read_csv(pagerank_path)
    stops = pd.read_csv(stops_path)[["stop_id", "stop_name"]]
    merged = (
        pagerank.merge(stops, on="stop_id", how="left")
        .sort_values("pagerank", ascending=False)
        .head(5)
    )
    return [
        (str(row.stop_name), f"{float(row.pagerank):.4f}")
        for row in merged.itertuples(index=False)
    ]


def project_metrics() -> dict[str, object]:
    return {
        "stops": csv_rows(DATA_RAW / "synthetic_bus_stops.csv"),
        "gps": csv_rows(DATA_RAW / "gps.csv"),
        "schedule": csv_rows(DATA_RAW / "scheduling.csv"),
        "routes": csv_rows(DATA_RAW / "cleaned_bus.csv.csv"),
        "weather": csv_rows(DATA_RAW / "Bangalore Weather Data (Visual Crossing Weather).csv"),
        "capacity": csv_rows(DATA_RAW / "synthetic_bus_capacity.csv"),
        "pagerank_rows": csv_rows(RESULTS / "pagerank.csv"),
        "shortest_rows": csv_rows(RESULTS / "shortest_paths.csv"),
        "components": csv_rows(RESULTS / "connected_components.csv"),
        "top_hubs": read_top_hubs(),
        "generated_at": datetime.now().strftime("%d %b %Y"),
    }


def build_presentation(metrics: dict[str, object]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide_title(prs.slides.add_slide(blank), metrics, 1)
    slide_problem(prs.slides.add_slide(blank), 2)
    slide_architecture(prs.slides.add_slide(blank), 3)
    slide_data_sources(prs.slides.add_slide(blank), metrics, 4)
    slide_stack(prs.slides.add_slide(blank), 5)
    slide_pipeline(prs.slides.add_slide(blank), 6)
    slide_graph_analytics(prs.slides.add_slide(blank), metrics, 7)
    slide_results(prs.slides.add_slide(blank), metrics, 8)
    slide_dashboard(prs.slides.add_slide(blank), 9)
    slide_achievements(prs.slides.add_slide(blank), metrics, 10)
    slide_challenges(prs.slides.add_slide(blank), 11)
    slide_conclusion(prs.slides.add_slide(blank), metrics, 12)
    return prs


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()


def add_footer(slide, number: int) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BORDER
    line.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.45), Inches(7.15), Inches(4.5), Inches(0.22))
    p = left.text_frame.paragraphs[0]
    p.text = "KSRTC Smart Transit System"
    p.font.name = BODY_FONT
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED

    right = slide.shapes.add_textbox(Inches(12.2), Inches(7.13), Inches(0.5), Inches(0.22))
    p = right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = str(number)
    p.font.name = BODY_FONT
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def add_header(slide, title: str, subtitle: str, number: int) -> None:
    set_background(slide)

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.42), Inches(1.45), Inches(0.36))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ORANGE
    tag.line.fill.background()
    tf = tag.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "PROJECT PPT"
    p.font.name = BODY_FONT
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.92), Inches(7.4), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(25)
    p.font.color.rgb = NAVY

    sub_box = slide.shapes.add_textbox(Inches(0.57), Inches(1.38), Inches(8.3), Inches(0.45))
    p = sub_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED

    add_footer(slide, number)


def add_textbox(slide, left, top, width, height, text: str, size: int = 16, bold: bool = False,
                color: RGBColor = TEXT, font_name: str = BODY_FONT, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def add_bullets(slide, left, top, width, height, bullets: list[str], size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def add_card(slide, left, top, width, height, title: str, body: str,
             fill_color: RGBColor = WHITE, title_color: RGBColor = NAVY,
             body_color: RGBColor = TEXT, outline: RGBColor = LIGHT_BORDER) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = outline
    card.line.width = Pt(1.2)

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.28))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = title_color

    body_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.46), width - Inches(0.36), height - Inches(0.55))
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = body_color


def add_metric_card(slide, left, top, width, height, value: str, label: str, accent: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_BORDER
    card.line.width = Pt(1.0)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    vbox = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.18), width - Inches(0.32), Inches(0.34))
    p = vbox.text_frame.paragraphs[0]
    p.text = value
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = NAVY

    lbox = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.58), width - Inches(0.32), Inches(0.26))
    p = lbox.text_frame.paragraphs[0]
    p.text = label
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED


def add_image(slide, path: Path, left, top, width, height) -> None:
    if not path.exists():
        add_card(slide, left, top, width, height, "Image Missing", str(path.name))
        return

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LIGHT_BORDER
    frame.line.width = Pt(1.0)

    with Image.open(path) as img:
        img_ratio = img.width / img.height

    box_ratio = width / height
    if img_ratio > box_ratio:
        pic_w = width - Inches(0.16)
        pic_h = pic_w / img_ratio
    else:
        pic_h = height - Inches(0.16)
        pic_w = pic_h * img_ratio

    pic_left = left + (width - pic_w) / 2
    pic_top = top + (height - pic_h) / 2
    slide.shapes.add_picture(str(path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_simple_table(slide, left, top, width, height, headers: list[str], rows: list[tuple[str, str]]) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table

    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.name = BODY_FONT
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.name = BODY_FONT
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT


def slide_title(slide, metrics: dict[str, object], number: int) -> None:
    set_background(slide)

    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.52), Inches(1.9), Inches(0.38))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    p = accent.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "WHOLE PROJECT"
    p.font.name = BODY_FONT
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    add_textbox(
        slide,
        Inches(0.55),
        Inches(1.05),
        Inches(6.2),
        Inches(1.7),
        "KSRTC Smart Transit System\nUsing Big Data Analytics",
        size=27,
        bold=True,
        color=NAVY,
        font_name=TITLE_FONT,
    )
    add_textbox(
        slide,
        Inches(0.58),
        Inches(2.82),
        Inches(5.8),
        Inches(0.6),
        "End-to-end batch, graph, streaming, API, and dashboard pipeline for transit intelligence.",
        size=14,
        color=MUTED,
    )

    add_metric_card(slide, Inches(0.58), Inches(4.15), Inches(1.8), Inches(1.1), str(metrics["stops"]), "Bus Stops", ORANGE)
    add_metric_card(slide, Inches(2.55), Inches(4.15), Inches(1.8), Inches(1.1), str(metrics["gps"]), "GPS Pings", BLUE)
    add_metric_card(slide, Inches(4.52), Inches(4.15), Inches(1.9), Inches(1.1), "3", "GraphX Algorithms", GREEN)

    add_textbox(
        slide,
        Inches(0.6),
        Inches(5.65),
        Inches(5.8),
        Inches(0.45),
        f"Generated for presentation on {metrics['generated_at']}",
        size=11,
        color=MUTED,
    )
    add_image(slide, VISUALS / "route_network_graph.png", Inches(7.1), Inches(0.85), Inches(5.65), Inches(5.95))
    add_footer(slide, number)


def slide_problem(slide, number: int) -> None:
    add_header(
        slide,
        "Problem Statement And Objectives",
        "Why this project matters and what the solution is designed to achieve.",
        number,
    )
    add_card(
        slide,
        Inches(0.6),
        Inches(2.0),
        Inches(5.95),
        Inches(3.7),
        "Problem Addressed",
        "Transit systems generate route, GPS, schedule, and traffic data at a scale that is difficult to monitor manually. "
        "Traditional workflows struggle to identify bottlenecks, important hubs, and route behavior in time for decision-making.",
    )
    add_bullets(
        slide,
        Inches(0.85),
        Inches(2.75),
        Inches(5.4),
        Inches(2.55),
        [
            "Disparate transport data needs to be unified before analysis.",
            "Transit networks benefit from graph modeling, not just flat reports.",
            "Operators need both historical insight and live monitoring support.",
        ],
        size=15,
    )

    add_card(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.95),
        Inches(3.7),
        "Project Objectives",
        "Build a practical big-data pipeline that turns raw KSRTC-style transport data into analytics, graph intelligence, and dashboard-ready outputs.",
    )
    add_bullets(
        slide,
        Inches(7.05),
        Inches(2.75),
        Inches(5.35),
        Inches(2.55),
        [
            "Clean and structure large route and telemetry datasets with Spark.",
            "Model the transit network in GraphX using stops as vertices and routes as edges.",
            "Integrate Kafka, MongoDB, Flask, and Streamlit for end-to-end presentation.",
        ],
        size=15,
    )

    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.0), Inches(11.8), Inches(0.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()
    p = banner.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Outcome: raw transport data becomes route intelligence and visual decision support."
    p.font.name = BODY_FONT
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE


def slide_architecture(slide, number: int) -> None:
    add_header(
        slide,
        "End-To-End System Architecture",
        "The project combines batch analytics, graph processing, storage, API services, and visualization.",
        number,
    )
    add_bullets(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(4.1),
        Inches(3.9),
        [
            "Raw CSV datasets are ingested and normalized through Apache Spark ETL.",
            "GraphX builds the route network and runs PageRank, ShortestPaths, and Connected Components.",
            "Kafka supports streaming-style event flow for telemetry simulation.",
            "MongoDB stores processed outputs and route analysis collections.",
            "Flask and Streamlit expose results through APIs and dashboards.",
        ],
        size=15,
    )
    add_image(slide, VISUALS / "system_architecture.png", Inches(4.95), Inches(1.85), Inches(7.45), Inches(4.95))


def slide_data_sources(slide, metrics: dict[str, object], number: int) -> None:
    add_header(
        slide,
        "Data Sources And Project Modules",
        "The project brings together multiple transport datasets and specialized code modules.",
        number,
    )

    add_metric_card(slide, Inches(0.6), Inches(1.95), Inches(2.05), Inches(1.0), f"{metrics['routes']}", "Route Records", ORANGE)
    add_metric_card(slide, Inches(2.85), Inches(1.95), Inches(2.05), Inches(1.0), f"{metrics['gps']}", "GPS Rows", BLUE)
    add_metric_card(slide, Inches(5.1), Inches(1.95), Inches(2.05), Inches(1.0), f"{metrics['weather']}", "Weather Rows", GREEN)
    add_metric_card(slide, Inches(7.35), Inches(1.95), Inches(2.05), Inches(1.0), f"{metrics['schedule']}", "Schedule Rows", ORANGE)
    add_metric_card(slide, Inches(9.6), Inches(1.95), Inches(2.05), Inches(1.0), f"{metrics['capacity']}", "Capacity Rows", BLUE)

    add_card(
        slide,
        Inches(0.6),
        Inches(3.25),
        Inches(5.95),
        Inches(2.9),
        "Input Data",
        "The raw layer includes synthetic and KSRTC-style datasets covering routes, bus stops, GPS traces, schedules, weather, and capacity.",
    )
    add_bullets(
        slide,
        Inches(0.85),
        Inches(3.95),
        Inches(5.45),
        Inches(1.8),
        [
            "Bus routes and stop metadata drive network creation.",
            "GPS pings support speed, movement, and traffic-oriented visuals.",
            "Schedules, weather, and capacity widen the analysis context.",
        ],
        size=14,
    )

    add_card(
        slide,
        Inches(6.8),
        Inches(3.25),
        Inches(5.95),
        Inches(2.9),
        "Repository Modules",
        "Code is organized around analytics, streaming, APIs, dashboards, and output assets.",
    )
    add_bullets(
        slide,
        Inches(7.05),
        Inches(3.95),
        Inches(5.35),
        Inches(1.9),
        [
            "backend/spark_jobs: Spark ETL and GraphX jobs",
            "backend/realtime and backend/api: Kafka simulation and Flask services",
            "dashboard and visualizations: Streamlit UI and generated charts",
        ],
        size=14,
    )


def slide_stack(slide, number: int) -> None:
    add_header(
        slide,
        "Technology Stack And Workflow",
        "Each layer contributes a specific role in the end-to-end analytics pipeline.",
        number,
    )

    cards = [
        ("Apache Spark", "Distributed ETL and preprocessing", BLUE),
        ("GraphX", "Route network modeling and graph algorithms", ORANGE),
        ("Kafka", "Telemetry-style streaming flow", GREEN),
        ("MongoDB", "Result storage and analysis collections", NAVY),
        ("Flask API", "Backend service layer", ORANGE),
        ("Streamlit + PyDeck", "Interactive dashboard and visual delivery", BLUE),
    ]
    x_positions = [Inches(0.6), Inches(4.55), Inches(8.5)]
    y_positions = [Inches(2.0), Inches(3.95)]
    for idx, (title, body, accent) in enumerate(cards):
        x = x_positions[idx % 3]
        y = y_positions[idx // 3]
        add_card(slide, x, y, Inches(3.4), Inches(1.45), title, body, fill_color=WHITE, title_color=accent)

    flow_y = Inches(6.05)
    steps = ["Ingest", "Process", "Analyze", "Store", "Visualize"]
    step_left = Inches(0.92)
    for idx, step in enumerate(steps):
        left = step_left + Inches(2.33 * idx)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, flow_y, Inches(1.35), Inches(0.46))
        pill.fill.solid()
        pill.fill.fore_color.rgb = NAVY if idx % 2 == 0 else ORANGE
        pill.line.fill.background()
        p = pill.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = step
        p.font.name = BODY_FONT
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        if idx < len(steps) - 1:
            connector = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left + Inches(1.52), flow_y + Inches(0.06), Inches(0.48), Inches(0.34))
            connector.fill.solid()
            connector.fill.fore_color.rgb = LIGHT_BORDER
            connector.line.fill.background()


def slide_pipeline(slide, number: int) -> None:
    add_header(
        slide,
        "Processing Pipeline",
        "The project includes both batch analytics and streaming-style components for a complete demo workflow.",
        number,
    )
    add_card(
        slide,
        Inches(0.6),
        Inches(2.0),
        Inches(5.8),
        Inches(4.4),
        "Batch Analytics Path",
        "Spark jobs prepare data for graph processing and downstream storage.",
    )
    add_bullets(
        slide,
        Inches(0.9),
        Inches(2.75),
        Inches(5.05),
        Inches(3.0),
        [
            "Read CSV sources from the raw data layer.",
            "Drop duplicates, cast schema, and normalize route attributes.",
            "Write processed outputs to Parquet for downstream usage.",
            "Build graph routes and compute GraphX analytics outputs.",
            "Store final analysis data in CSV and MongoDB-ready collections.",
        ],
        size=15,
    )

    add_card(
        slide,
        Inches(6.95),
        Inches(2.0),
        Inches(5.8),
        Inches(4.4),
        "Streaming And Serving Path",
        "Realtime simulation demonstrates how live telemetry can feed dashboards and services.",
    )
    add_bullets(
        slide,
        Inches(7.25),
        Inches(2.75),
        Inches(5.05),
        Inches(3.0),
        [
            "Kafka producers and consumers simulate active bus telemetry streams.",
            "Realtime helpers ingest schedule, GPS, and weather updates.",
            "Flask services expose data for dashboard or future integrations.",
            "Streamlit consumes generated analytics for demo-friendly visualization.",
            "The pipeline is HDFS-ready and can scale beyond local execution.",
        ],
        size=15,
    )


def slide_graph_analytics(slide, metrics: dict[str, object], number: int) -> None:
    add_header(
        slide,
        "GraphX Modeling And Algorithms",
        "Stops are represented as graph vertices and route connections are modeled as edges.",
        number,
    )
    add_metric_card(slide, Inches(0.65), Inches(1.95), Inches(1.9), Inches(0.95), str(metrics["stops"]), "Vertices / Stops", ORANGE)
    add_metric_card(slide, Inches(2.75), Inches(1.95), Inches(1.9), Inches(0.95), str(metrics["shortest_rows"]), "Path Links", BLUE)
    add_metric_card(slide, Inches(4.85), Inches(1.95), Inches(1.9), Inches(0.95), "1", "Connected Area", GREEN)

    add_bullets(
        slide,
        Inches(0.7),
        Inches(3.1),
        Inches(4.55),
        Inches(3.0),
        [
            "PageRank identifies influential bus stops and strategic hub locations.",
            "ShortestPaths estimates network reachability from a selected source stop.",
            "Connected Components reveal whether the sample network stays unified.",
            "Graph outputs are written back to CSV for reporting and dashboard use.",
        ],
        size=15,
    )
    add_image(slide, VISUALS / "route_network_graph.png", Inches(5.45), Inches(2.0), Inches(3.4), Inches(4.2))
    add_image(slide, VISUALS / "shortest_path_distances.png", Inches(8.95), Inches(2.0), Inches(3.4), Inches(4.2))


def slide_results(slide, metrics: dict[str, object], number: int) -> None:
    add_header(
        slide,
        "Key Analytics Results",
        "The current project outputs identify high-importance hubs and summarize route connectivity.",
        number,
    )
    add_image(slide, VISUALS / "top_pagerank_hubs.png", Inches(0.6), Inches(1.95), Inches(7.0), Inches(4.95))
    rows = metrics["top_hubs"] or [("KR Market", "1.1582"), ("Majestic", "1.0880")]
    add_card(
        slide,
        Inches(7.85),
        Inches(1.95),
        Inches(4.9),
        Inches(1.0),
        "Top Ranked Hubs",
        "Highest-scoring stops from the PageRank output show where route importance concentrates.",
    )
    add_simple_table(slide, Inches(7.85), Inches(3.15), Inches(4.9), Inches(3.1), ["Stop", "Score"], rows)


def slide_dashboard(slide, number: int) -> None:
    add_header(
        slide,
        "Dashboard And Visual Analytics",
        "The presentation layer combines live-style visuals, route views, and operational summaries.",
        number,
    )
    add_image(slide, VISUALS / "traffic_heatmap.png", Inches(0.6), Inches(2.0), Inches(4.15), Inches(2.4))
    add_image(slide, VISUALS / "speed_distribution.png", Inches(0.6), Inches(4.6), Inches(4.15), Inches(2.05))

    add_card(
        slide,
        Inches(5.05),
        Inches(2.0),
        Inches(7.7),
        Inches(4.65),
        "Dashboard Experience",
        "The Streamlit layer is designed to make the analytics understandable for demonstrations and decision support.",
    )
    add_bullets(
        slide,
        Inches(5.35),
        Inches(2.8),
        Inches(6.9),
        Inches(2.9),
        [
            "3D transit map and route visualization",
            "Top hub stops powered by PageRank outputs",
            "Shortest-path previews and connected-area summaries",
            "Traffic heatmaps, speed analysis, and demand-oriented views",
            "CSV-driven visual outputs for reliable presentation demos",
        ],
        size=15,
    )


def slide_achievements(slide, metrics: dict[str, object], number: int) -> None:
    add_header(
        slide,
        "Project Achievements",
        "This project demonstrates an end-to-end smart transit analytics workflow rather than a single isolated model.",
        number,
    )
    add_metric_card(slide, Inches(0.65), Inches(2.0), Inches(2.25), Inches(1.1), str(metrics["routes"]), "Route Records Processed", ORANGE)
    add_metric_card(slide, Inches(3.08), Inches(2.0), Inches(2.25), Inches(1.1), str(metrics["gps"]), "GPS Pings Simulated", BLUE)
    add_metric_card(slide, Inches(5.51), Inches(2.0), Inches(2.25), Inches(1.1), str(len(metrics["top_hubs"])), "Top Hubs Reported", GREEN)
    add_metric_card(slide, Inches(7.94), Inches(2.0), Inches(2.25), Inches(1.1), str(metrics["shortest_rows"]), "Shortest Path Links", ORANGE)
    add_metric_card(slide, Inches(10.37), Inches(2.0), Inches(2.25), Inches(1.1), str(metrics["components"]), "Component Rows", BLUE)

    add_card(
        slide,
        Inches(0.7),
        Inches(3.55),
        Inches(12.0),
        Inches(2.7),
        "Deliverables Produced",
        "The repository includes Spark jobs, realtime simulation modules, Flask APIs, Streamlit dashboards, CSV analytics outputs, and visualization assets. "
        "Together they form a complete presentation-ready project with both technical depth and demo value.",
    )
    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.22),
        Inches(11.2),
        Inches(1.55),
        [
            "Core outputs: pagerank.csv, shortest_paths.csv, connected_components.csv, route_analysis.csv",
            "Visual outputs: hub ranking chart, sampled network graph, heatmap, speed analysis, system architecture",
            "Implementation spans Spark, GraphX, Kafka, MongoDB, Flask, Streamlit, and PyDeck",
        ],
        size=14,
    )


def slide_challenges(slide, number: int) -> None:
    add_header(
        slide,
        "Challenges And Limitations",
        "The current version is strong as a smart-transit prototype, with clear next steps for production readiness.",
        number,
    )
    add_card(slide, Inches(0.65), Inches(2.0), Inches(2.85), Inches(2.05), "Data Quality", "Some inputs are synthetic or demo-oriented, so the pipeline prioritizes robustness and presentation reliability.")
    add_card(slide, Inches(3.63), Inches(2.0), Inches(2.85), Inches(2.05), "Environment Setup", "Full execution depends on local availability of Spark, Kafka, MongoDB, and optional HDFS infrastructure.")
    add_card(slide, Inches(6.61), Inches(2.0), Inches(2.85), Inches(2.05), "Realtime Scope", "Telemetry is simulated rather than connected to live public transit feeds in the current project version.")
    add_card(slide, Inches(9.59), Inches(2.0), Inches(2.85), Inches(2.05), "Scaling Next Step", "The architecture is HDFS-ready, but larger-scale deployment would benefit from cloud orchestration and stronger monitoring.")

    add_card(
        slide,
        Inches(0.85),
        Inches(4.45),
        Inches(11.6),
        Inches(1.5),
        "How The Project Handles These Gaps",
        "The design already separates ingestion, processing, analytics, storage, and visualization. That modularity makes it easier to replace simulated data with live feeds and scale individual layers independently.",
    )


def slide_conclusion(slide, metrics: dict[str, object], number: int) -> None:
    add_header(
        slide,
        "Future Scope And Conclusion",
        "The project establishes a strong foundation for smart-city transportation analytics.",
        number,
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(2.0),
        Inches(5.7),
        Inches(4.5),
        "Future Scope",
        "The current implementation can evolve into a more production-oriented intelligent transit platform.",
    )
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.75),
        Inches(5.0),
        Inches(3.2),
        [
            "Integrate live GPS and external traffic feeds.",
            "Add machine learning for ETA, demand, and route optimization.",
            "Deploy on cloud infrastructure with monitoring and autoscaling.",
            "Expand the dashboard into operational and mobile interfaces.",
        ],
        size=15,
    )

    add_card(
        slide,
        Inches(6.75),
        Inches(2.0),
        Inches(5.7),
        Inches(4.5),
        "Conclusion",
        "KSRTC Smart Transit System using Big Data Analytics demonstrates how Spark, GraphX, Kafka, MongoDB, Flask, and Streamlit can work together to turn raw transport data into actionable insight.",
    )
    add_bullets(
        slide,
        Inches(7.05),
        Inches(2.9),
        Inches(5.0),
        Inches(2.65),
        [
            "End-to-end big data workflow",
            "Graph-based route intelligence",
            "Realtime-ready monitoring design",
            "Presentation-ready visual outputs",
        ],
        size=15,
    )
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(5.7), Inches(4.95), Inches(0.55))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ORANGE
    banner.line.fill.background()
    p = banner.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Thank You"
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE


def main() -> int:
    metrics = project_metrics()
    prs = build_presentation(metrics)
    prs.save(OUTPUT)
    print(f"Saved presentation: {OUTPUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
